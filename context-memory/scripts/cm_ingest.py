#!/usr/bin/env python3
"""
cm_ingest.py — Extract knowledge from documents and store as Memory Nodes.

Supports: PDF, DOCX, XLSX, PPTX, CSV/TSV, JSON, TXT/MD, images (via description).

Usage:
  python3 cm_ingest.py --file report.pdf [--type architecture] [--tags "spec,api"]
  python3 cm_ingest.py --file data.xlsx --type config --tags "schema,daten"
  python3 cm_ingest.py --file slides.pptx --split slides
  python3 cm_ingest.py --file requirements.docx --split sections
"""

import argparse
import json
import os
import re
import subprocess
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))
from cm_core import (
    get_workspace, ensure_workspace, generate_id, load_index, save_index,
    load_tree, save_tree, save_node, add_history_entry,
    normalize_tags, content_hash, find_similar_nodes,
    VALID_TYPES, _now
)


# ─── File type detection ─────────────────────────────────────────────

def detect_file_type(filepath: Path) -> str:
    """Detect file type from extension."""
    ext = filepath.suffix.lower()
    type_map = {
        ".pdf": "pdf",
        ".docx": "docx", ".doc": "doc",
        ".xlsx": "xlsx", ".xlsm": "xlsx", ".xls": "xls", ".ods": "ods",
        ".pptx": "pptx", ".ppt": "ppt",
        ".csv": "csv", ".tsv": "tsv",
        ".json": "json", ".jsonl": "jsonl",
        ".txt": "text", ".md": "text", ".log": "text", ".rst": "text",
        ".jpg": "image", ".jpeg": "image", ".png": "image",
        ".gif": "image", ".webp": "image", ".bmp": "image",
        ".epub": "epub", ".odt": "odt", ".rtf": "rtf",
        ".html": "html", ".htm": "html",
        ".xml": "xml", ".yaml": "yaml", ".yml": "yaml", ".toml": "toml",
        ".py": "code", ".js": "code", ".ts": "code", ".java": "code",
        ".go": "code", ".rs": "code", ".rb": "code", ".swift": "code",
        ".kt": "code", ".cs": "code", ".cpp": "code", ".c": "code",
        ".sh": "code", ".bash": "code",
    }
    return type_map.get(ext, "unknown")


# ─── Extractors ──────────────────────────────────────────────────────

def extract_pdf(filepath: Path, max_pages: int = 50) -> dict:
    """Extract text from PDF."""
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(filepath))
        total_pages = len(reader.pages)
        
        pages = []
        for i, page in enumerate(reader.pages[:max_pages]):
            text = page.extract_text()
            if text and text.strip():
                pages.append({"page": i + 1, "text": text.strip()})
        
        # If pypdf extracted nothing, try pdftotext
        if not pages:
            result = subprocess.run(
                ["pdftotext", "-layout", str(filepath), "-"],
                capture_output=True, text=True, timeout=30
            )
            if result.stdout.strip():
                pages = [{"page": 1, "text": result.stdout.strip()[:50000]}]
        
        return {
            "source": filepath.name,
            "type": "pdf",
            "total_pages": total_pages,
            "extracted_pages": len(pages),
            "pages": pages,
            "full_text": "\n\n".join(p["text"] for p in pages)
        }
    except Exception as e:
        return {"error": f"PDF extraction failed: {e}", "source": filepath.name}


def extract_docx(filepath: Path) -> dict:
    """Extract text from DOCX, preserving section structure."""
    try:
        # Use pandoc for best markdown conversion
        result = subprocess.run(
            ["pandoc", str(filepath), "-t", "markdown", "--wrap=none"],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0 and result.stdout.strip():
            text = result.stdout.strip()
            sections = _split_markdown_sections(text)
            return {
                "source": filepath.name,
                "type": "docx",
                "sections": sections,
                "full_text": text
            }
        
        # Fallback: python-docx
        from docx import Document
        doc = Document(str(filepath))
        paragraphs = []
        current_section = {"heading": "Document", "content": []}
        sections = []
        
        for para in doc.paragraphs:
            if para.style.name.startswith("Heading"):
                if current_section["content"]:
                    sections.append(current_section)
                current_section = {"heading": para.text, "content": []}
            elif para.text.strip():
                current_section["content"].append(para.text.strip())
        
        if current_section["content"]:
            sections.append(current_section)
        
        for s in sections:
            s["content"] = "\n".join(s["content"])
        
        full_text = "\n\n".join(
            f"## {s['heading']}\n{s['content']}" for s in sections
        )
        
        return {
            "source": filepath.name,
            "type": "docx",
            "sections": sections,
            "full_text": full_text
        }
    except Exception as e:
        return {"error": f"DOCX extraction failed: {e}", "source": filepath.name}


def extract_xlsx(filepath: Path, max_rows: int = 100) -> dict:
    """Extract structure and data from Excel files."""
    try:
        import pandas as pd
        ext = filepath.suffix.lower()
        
        engine = None
        if ext == ".xls":
            engine = "xlrd"
        elif ext == ".ods":
            engine = "odf"
        
        kwargs = {"engine": engine} if engine else {}
        xls = pd.ExcelFile(str(filepath), **kwargs)
        
        sheets_info = []
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name, nrows=max_rows)
            sheets_info.append({
                "name": sheet_name,
                "columns": list(df.columns),
                "dtypes": {col: str(dtype) for col, dtype in df.dtypes.items()},
                "shape": list(df.shape),
                "sample": df.head(5).to_string(),
                "nulls": {col: int(v) for col, v in df.isnull().sum().items() if v > 0}
            })
        
        full_text = ""
        for s in sheets_info:
            full_text += (
                f"## Sheet: {s['name']}\n"
                f"Shape: {s['shape'][0]} rows × {s['shape'][1]} columns\n"
                f"Columns: {', '.join(s['columns'])}\n"
                f"Types: {json.dumps(s['dtypes'])}\n"
                f"{'Null values: ' + json.dumps(s['nulls']) if s['nulls'] else 'No null values'}\n\n"
                f"Sample:\n```\n{s['sample']}\n```\n\n"
            )
        
        return {
            "source": filepath.name,
            "type": "xlsx",
            "sheets": sheets_info,
            "full_text": full_text
        }
    except Exception as e:
        return {"error": f"Excel extraction failed: {e}", "source": filepath.name}


def extract_pptx(filepath: Path) -> dict:
    """Extract text from PowerPoint slides."""
    try:
        from pptx import Presentation
        prs = Presentation(str(filepath))
        
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            notes = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
            
            if texts or notes:
                slides.append({
                    "slide": i + 1,
                    "content": texts,
                    "notes": notes
                })
        
        full_text = ""
        for s in slides:
            full_text += f"## Slide {s['slide']}\n"
            full_text += "\n".join(s["content"])
            if s["notes"]:
                full_text += f"\n\n*Speaker Notes:* {s['notes']}"
            full_text += "\n\n"
        
        return {
            "source": filepath.name,
            "type": "pptx",
            "total_slides": len(prs.slides),
            "slides": slides,
            "full_text": full_text
        }
    except Exception as e:
        return {"error": f"PPTX extraction failed: {e}", "source": filepath.name}


def extract_csv(filepath: Path, max_rows: int = 200) -> dict:
    """Extract structure and data from CSV/TSV."""
    try:
        import pandas as pd
        sep = "\t" if filepath.suffix.lower() == ".tsv" else ","
        df = pd.read_csv(str(filepath), sep=sep, nrows=max_rows)
        
        full_text = (
            f"**File:** {filepath.name}\n"
            f"**Shape:** {df.shape[0]} rows × {df.shape[1]} columns\n"
            f"**Columns:** {', '.join(df.columns)}\n"
            f"**Types:**\n{df.dtypes.to_string()}\n\n"
            f"**Statistics:**\n{df.describe().to_string()}\n\n"
            f"**First 5 rows:**\n```\n{df.head().to_string()}\n```"
        )
        
        return {
            "source": filepath.name,
            "type": "csv",
            "columns": list(df.columns),
            "shape": list(df.shape),
            "full_text": full_text
        }
    except Exception as e:
        return {"error": f"CSV extraction failed: {e}", "source": filepath.name}


def extract_json_file(filepath: Path) -> dict:
    """Extract structure from JSON."""
    try:
        content = filepath.read_text(encoding="utf-8")
        data = json.loads(content)
        
        def describe_structure(obj, depth=0, max_depth=3):
            if depth > max_depth:
                return "..."
            if isinstance(obj, dict):
                items = {k: describe_structure(v, depth + 1) for k, v in list(obj.items())[:20]}
                return items
            elif isinstance(obj, list):
                if obj:
                    return [describe_structure(obj[0], depth + 1), f"... ({len(obj)} items)"]
                return "[]"
            else:
                return type(obj).__name__
        
        structure = describe_structure(data)
        
        full_text = (
            f"**File:** {filepath.name}\n"
            f"**Type:** {type(data).__name__}\n"
            f"**Structure:**\n```json\n{json.dumps(structure, indent=2, ensure_ascii=False)[:3000]}\n```\n\n"
        )
        
        if isinstance(data, dict):
            full_text += f"**Top-level keys:** {', '.join(list(data.keys())[:30])}"
        elif isinstance(data, list):
            full_text += f"**Array length:** {len(data)}"
        
        return {
            "source": filepath.name,
            "type": "json",
            "structure": structure,
            "full_text": full_text
        }
    except Exception as e:
        return {"error": f"JSON extraction failed: {e}", "source": filepath.name}


def extract_text(filepath: Path, max_chars: int = 50000) -> dict:
    """Extract from plain text / markdown / code files."""
    try:
        content = filepath.read_text(encoding="utf-8")[:max_chars]
        ext = filepath.suffix.lower()
        
        result = {
            "source": filepath.name,
            "type": "text",
            "full_text": content,
            "lines": content.count("\n") + 1,
            "chars": len(content)
        }
        
        # For markdown, try to extract sections
        if ext in (".md", ".rst"):
            result["sections"] = _split_markdown_sections(content)
        
        return result
    except UnicodeDecodeError:
        return {"error": "Binary file, cannot read as text", "source": filepath.name}
    except Exception as e:
        return {"error": f"Text extraction failed: {e}", "source": filepath.name}


def extract_epub(filepath: Path) -> dict:
    """Extract text from EPUB."""
    try:
        result = subprocess.run(
            ["pandoc", str(filepath), "-t", "plain", "--wrap=none"],
            capture_output=True, text=True, timeout=60
        )
        if result.returncode == 0:
            text = result.stdout.strip()[:50000]
            return {
                "source": filepath.name,
                "type": "epub",
                "full_text": text,
                "chars": len(text)
            }
        return {"error": f"pandoc failed: {result.stderr}", "source": filepath.name}
    except Exception as e:
        return {"error": f"EPUB extraction failed: {e}", "source": filepath.name}


def extract_html(filepath: Path) -> dict:
    """Extract text from HTML."""
    try:
        result = subprocess.run(
            ["pandoc", str(filepath), "-f", "html", "-t", "markdown", "--wrap=none"],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0:
            text = result.stdout.strip()[:50000]
            return {
                "source": filepath.name,
                "type": "html",
                "sections": _split_markdown_sections(text),
                "full_text": text
            }
        # Fallback: basic regex strip
        content = filepath.read_text(encoding="utf-8")
        text = re.sub(r'<[^>]+>', ' ', content)
        text = re.sub(r'\s+', ' ', text).strip()[:50000]
        return {"source": filepath.name, "type": "html", "full_text": text}
    except Exception as e:
        return {"error": f"HTML extraction failed: {e}", "source": filepath.name}


def extract_config(filepath: Path) -> dict:
    """Extract from YAML/TOML config files."""
    try:
        content = filepath.read_text(encoding="utf-8")[:20000]
        return {
            "source": filepath.name,
            "type": "config",
            "full_text": f"```{filepath.suffix.lstrip('.')}\n{content}\n```"
        }
    except Exception as e:
        return {"error": f"Config extraction failed: {e}", "source": filepath.name}


def extract_image(filepath: Path) -> dict:
    """Get image metadata (actual vision analysis requires Claude)."""
    try:
        from PIL import Image
        img = Image.open(str(filepath))
        return {
            "source": filepath.name,
            "type": "image",
            "dimensions": f"{img.size[0]}×{img.size[1]}",
            "mode": img.mode,
            "format": img.format,
            "full_text": (
                f"**Image:** {filepath.name}\n"
                f"**Dimensions:** {img.size[0]}×{img.size[1]} px\n"
                f"**Mode:** {img.mode}\n"
                f"**Format:** {img.format}\n\n"
                f"*Note: For content description, let Claude analyze the image directly.*"
            )
        }
    except ImportError:
        return {
            "source": filepath.name,
            "type": "image",
            "full_text": f"**Image:** {filepath.name}\n*PIL not available for metadata. Let Claude analyze directly.*"
        }
    except Exception as e:
        return {"error": f"Image metadata failed: {e}", "source": filepath.name}


# ─── Helpers ─────────────────────────────────────────────────────────

def _split_markdown_sections(text: str) -> list:
    """Split markdown text into sections by headings."""
    sections = []
    current = {"heading": "Introduction", "content": []}
    
    for line in text.split("\n"):
        heading_match = re.match(r'^(#{1,4})\s+(.+)$', line)
        if heading_match:
            if current["content"]:
                current["content"] = "\n".join(current["content"]).strip()
                sections.append(current)
            current = {"heading": heading_match.group(2).strip(), "content": []}
        else:
            current["content"].append(line)
    
    if current["content"]:
        current["content"] = "\n".join(current["content"]).strip()
        sections.append(current)
    
    return [s for s in sections if s["content"]]


def _truncate(text: str, max_chars: int = 10000) -> str:
    """Truncate text to max chars with indicator."""
    if len(text) <= max_chars:
        return text
    return text[:max_chars] + f"\n\n... (truncated, {len(text) - max_chars} chars omitted)"


# ─── Main dispatch ───────────────────────────────────────────────────

EXTRACTORS = {
    "pdf": extract_pdf,
    "docx": extract_docx,
    "doc": extract_docx,  # pandoc handles .doc too in many cases
    "xlsx": extract_xlsx,
    "xls": extract_xlsx,
    "ods": extract_xlsx,
    "pptx": extract_pptx,
    "csv": extract_csv,
    "tsv": extract_csv,
    "json": extract_json_file,
    "jsonl": extract_json_file,
    "text": extract_text,
    "code": extract_text,
    "epub": extract_epub,
    "odt": extract_epub,  # pandoc handles odt
    "html": extract_html,
    "xml": extract_config,
    "yaml": extract_config,
    "toml": extract_config,
    "rtf": extract_epub,  # pandoc handles rtf
    "image": extract_image,
}


def main():
    parser = argparse.ArgumentParser(description="Ingest documents into Context Memory")
    parser.add_argument("--file", "-f", required=True, help="Path to file")
    parser.add_argument("--type", "-t", default=None, choices=VALID_TYPES,
                        help="Node type (auto-detected if omitted)")
    parser.add_argument("--tags", default="", help="Comma-separated tags")
    parser.add_argument("--relevance", "-r", default="medium",
                        choices=["critical", "high", "medium", "low"])
    parser.add_argument("--title", default=None, help="Custom title (default: filename)")
    parser.add_argument("--split", "-s", default="single",
                        choices=["single", "sections", "slides", "sheets", "pages"],
                        help="How to split content into nodes")
    parser.add_argument("--max-nodes", type=int, default=20,
                        help="Max nodes to create when splitting")
    parser.add_argument("--project-name", "-n", default=None)
    parser.add_argument("--path", "-p", default=None)
    parser.add_argument("--dry-run", action="store_true")
    parser.add_argument("--force", action="store_true", help="Skip duplicate check")
    parser.add_argument("--description", "-d", default=None,
                        help="For images: description of what the image shows")
    args = parser.parse_args()

    filepath = Path(args.file).resolve()
    if not filepath.exists():
        print(f"❌ File not found: {filepath}")
        sys.exit(1)

    file_type = detect_file_type(filepath)
    if file_type == "unknown":
        print(f"⚠️  Unknown file type: {filepath.suffix}")
        print(f"   Trying as plain text...")
        file_type = "text"

    extractor = EXTRACTORS.get(file_type)
    if not extractor:
        print(f"❌ No extractor for type '{file_type}'")
        sys.exit(1)

    # File info
    size = filepath.stat().st_size
    size_str = f"{size / 1024:.1f} KB" if size < 1024 * 1024 else f"{size / 1024 / 1024:.1f} MB"
    print(f"📄 Ingesting: {filepath.name} ({size_str}, type: {file_type})")

    # Extract content
    extracted = extractor(filepath)
    
    if "error" in extracted:
        print(f"❌ {extracted['error']}")
        sys.exit(1)

    # For images with user-provided description
    if file_type == "image" and args.description:
        extracted["full_text"] += f"\n\n**Content Description:** {args.description}"

    # Auto-detect node type if not specified
    node_type = args.type
    if not node_type:
        type_hints = {
            "pdf": "architecture",
            "docx": "architecture",
            "xlsx": "config",
            "pptx": "architecture",
            "csv": "config",
            "json": "config",
            "code": "pattern",
            "image": "architecture",
        }
        node_type = type_hints.get(file_type, "architecture")

    # Prepare workspace
    ws = get_workspace(args.project_name, args.path)
    ensure_workspace(ws)

    # Build auto-tags from filename and file type
    auto_tags = [file_type, filepath.suffix.lstrip(".")]
    name_parts = re.sub(r'[^\w]', ' ', filepath.stem.lower()).split()
    auto_tags.extend(p for p in name_parts if len(p) > 2)
    all_tags = normalize_tags(args.tags + "," + ",".join(auto_tags))

    # Decide splitting strategy
    nodes_to_create = []
    base_title = args.title or filepath.stem.replace("-", " ").replace("_", " ").title()

    if args.split == "single" or not _has_parts(extracted, args.split):
        # Single node with full content
        nodes_to_create.append({
            "title": base_title,
            "content": _truncate(extracted.get("full_text", "No content extracted")),
            "tags": all_tags
        })

    elif args.split == "sections" and extracted.get("sections"):
        for i, section in enumerate(extracted["sections"][:args.max_nodes]):
            section_content = section.get("content", "")
            if isinstance(section_content, list):
                section_content = "\n".join(section_content)
            if not section_content.strip():
                continue
            nodes_to_create.append({
                "title": f"{base_title} — {section.get('heading', f'Section {i+1}')}",
                "content": _truncate(section_content),
                "tags": all_tags + normalize_tags(section.get("heading", ""))
            })

    elif args.split == "slides" and extracted.get("slides"):
        for slide in extracted["slides"][:args.max_nodes]:
            content_parts = slide.get("content", [])
            if isinstance(content_parts, list):
                content = "\n".join(content_parts)
            else:
                content = str(content_parts)
            notes = slide.get("notes", "")
            if notes:
                content += f"\n\n*Speaker Notes:* {notes}"
            if not content.strip():
                continue
            nodes_to_create.append({
                "title": f"{base_title} — Slide {slide.get('slide', '?')}",
                "content": content,
                "tags": all_tags
            })

    elif args.split == "sheets" and extracted.get("sheets"):
        for sheet in extracted["sheets"][:args.max_nodes]:
            content = (
                f"**Columns:** {', '.join(sheet.get('columns', []))}\n"
                f"**Shape:** {sheet.get('shape', '?')}\n"
                f"**Types:** {json.dumps(sheet.get('dtypes', {}))}\n\n"
                f"**Sample:**\n```\n{sheet.get('sample', '')}\n```"
            )
            nodes_to_create.append({
                "title": f"{base_title} — Sheet: {sheet.get('name', '?')}",
                "content": content,
                "tags": all_tags + [sheet.get("name", "").lower().replace(" ", "-")]
            })

    elif args.split == "pages" and extracted.get("pages"):
        for page in extracted["pages"][:args.max_nodes]:
            if not page.get("text", "").strip():
                continue
            nodes_to_create.append({
                "title": f"{base_title} — Page {page.get('page', '?')}",
                "content": _truncate(page["text"]),
                "tags": all_tags
            })

    if not nodes_to_create:
        print("📭 No content extracted. The file might be empty or in an unsupported format.")
        return

    # Dry run
    if args.dry_run:
        print(f"\n🔎 Would create {len(nodes_to_create)} node(s):\n")
        for nd in nodes_to_create:
            preview = nd["content"][:120].replace("\n", " ")
            print(f"   [{node_type}] {nd['title']}")
            print(f"   Tags: {', '.join(nd['tags'])}")
            print(f"   Preview: {preview}...")
            print()
        return

    # Create nodes
    index = load_index(ws)
    tree = load_tree(ws)
    added = []

    for nd in nodes_to_create:
        # Duplicate check
        if not args.force:
            similar = find_similar_nodes(ws, nd["title"], nd["content"])
            if similar:
                print(f"   ⚠️  Skipping '{nd['title']}' — similar to {similar[0]['id']}")
                continue

        node_id = generate_id(node_type, ws, index)
        tags = normalize_tags(",".join(nd["tags"]))
        
        node_md = (
            f"# {nd['title']}\n\n"
            f"**Type:** {node_type} | **Relevance:** {args.relevance} | **Status:** active\n"
            f"**Tags:** {', '.join(tags)}\n"
            f"**Created:** {_now()} | **Source:** ingested from {filepath.name}\n\n"
            f"---\n\n{nd['content']}\n"
        )
        save_node(ws, node_id, node_md)
        
        index.setdefault("nodes", {})[node_id] = {
            "title": nd["title"],
            "type": node_type,
            "tags": tags,
            "relevance": args.relevance,
            "status": "active",
            "created": _now(),
            "updated": _now(),
            "parent": None,
            "hash": content_hash(nd["content"]),
            "source": f"ingested:{filepath.name}"
        }
        
        tree.setdefault("root_children", []).append({"id": node_id, "children": []})
        added.append((node_id, nd["title"]))

    save_index(ws, index)
    save_tree(ws, tree)

    for nid, title in added:
        add_history_entry(ws, "ingest", nid, f"Ingested from {filepath.name}: {title}")

    print(f"\n✅ Ingested {len(added)} node(s) from {filepath.name}:\n")
    for nid, title in added:
        print(f"   📝 {nid}: {title}")
    
    if len(nodes_to_create) > len(added):
        skipped = len(nodes_to_create) - len(added)
        print(f"\n   ⚠️  {skipped} node(s) skipped (duplicates)")
    
    print(f"\n   Workspace: {ws}")


def _has_parts(extracted: dict, split_mode: str) -> bool:
    """Check if the extracted data has parts for the given split mode."""
    part_keys = {
        "sections": "sections",
        "slides": "slides",
        "sheets": "sheets",
        "pages": "pages"
    }
    key = part_keys.get(split_mode)
    if key:
        parts = extracted.get(key, [])
        return len(parts) > 1
    return False


if __name__ == "__main__":
    main()
